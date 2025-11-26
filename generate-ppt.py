from pptx import Presentation

ppt_path_full = "/mnt/data/python_cloud_full_summary.pptx"

prs = Presentation()

# Slide titles and content
slides_content = [
    ("Python for Amazon Web Services", 
     "• Python integrates with AWS through the boto3 SDK\n"
     "• Supports automation for EC2, S3, Lambda, DynamoDB, IAM, SNS, SQS, CloudWatch\n"
     "• Enables provisioning, configuration, monitoring, and serverless workflows\n"
     "• Common use cases: Upload/download objects to S3, Start/stop EC2 instances, Trigger Lambda functions, Publish messages to SNS, Push/receive messages in SQS\n"
     "• Security: IAM user or role with required permissions; credentials via environment variables or IAM roles\n"
     "• Execution environments: Local machine, Lambda runtime, EC2 or containers"
    ),
    ("Python for Google Cloud",
     "• Python integrates with Google Cloud through google-cloud-* client libraries\n"
     "• Supports Compute Engine, Cloud Storage, BigQuery, Pub/Sub, Cloud Functions\n"
     "• Common operations: Upload/download objects to Cloud Storage, Run queries on BigQuery, Publish/subscribe to Pub/Sub topics, Invoke Cloud Functions\n"
     "• Authentication: service account key or Application Default Credentials\n"
     "• Execution environments: Local machine, Cloud Functions runtime, Compute Engine or containers"
    ),
    ("Python for Windows Azure",
     "• Python integrates with Azure via the Azure SDK for Python (azure-*)\n"
     "• Supports Azure Resource Manager, VMs, Storage Accounts, Blob Storage, Functions, Service Bus, Cosmos DB\n"
     "• Common operations: Upload/download blobs, Manage VMs, Send/receive messages, Insert/query documents in Cosmos DB\n"
     "• Authentication options: Service Principal, Managed Identity, Azure CLI authentication\n"
     "• Execution environments: Local machine, Azure Functions, Azure VM, App Service, Container Apps"
    ),
    ("Python for MapReduce",
     "• Python supports MapReduce using Hadoop Streaming, PySpark, and MRJob\n"
     "• Enables distributed processing across clusters for large datasets\n"
     "• Stages: Map (transform input to key-value pairs), Shuffle (group by key), Reduce (aggregate values)\n"
     "• Execution environments: Hadoop cluster, AWS EMR, Google Dataproc, local PySpark\n"
     "• Use cases: Log processing, Word frequency analysis, Aggregation of large datasets\n"
     "• Libraries: pyspark, mrjob, Hadoop Streaming"
    ),
    ("Python Packages of Interest",
     "• Core categories:\n"
     "  - Data Analysis: pandas, numpy\n"
     "  - Web Requests & APIs: requests, httpx\n"
     "  - Cloud SDKs: boto3, google-cloud-*, azure-*\n"
     "  - Machine Learning: scikit-learn, tensorflow, torch\n"
     "  - Web Frameworks: django, flask, fastapi\n"
     "  - Task Automation: fabric, invoke\n"
     "  - Messaging & Queues: kafka-python, pika\n"
     "  - Databases: sqlalchemy, pymongo, psycopg2\n"
     "• Widely used in enterprise systems, DevOps, data engineering, backend development"
    ),
    ("Python Web Application Framework – Django",
     "• High-level Python framework for rapid web development\n"
     "• Core components: Models (ORM), Views, Templates, URL routing, Forms, Admin site\n"
     "• Built-in security: CSRF, XSS protection, SQL injection defense\n"
     "• Supports PostgreSQL, MySQL, SQLite, Oracle\n"
     "• Suitable for enterprise systems, APIs, dashboards, CMS, e-commerce"
    ),
    ("Development with Django",
     "• Typical workflow: Create project → Create app → Define models → run migrations → Implement views → Configure URLs → Build templates → Use ORM → Admin configuration → Deploy\n"
     "• Recommended project structure:\n"
     "  project/\n"
     "      settings.py\n"
     "      urls.py\n"
     "      app/\n"
     "          models.py\n"
     "          views.py\n"
     "          urls.py\n"
     "          templates/"
    )
]

# Add slides
for title, content in slides_content:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

prs.save(ppt_path_full)
ppt_path_full
